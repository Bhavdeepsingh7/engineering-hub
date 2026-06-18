from pptx import Presentation


def extract_pptx_text(file_path: str):
    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text